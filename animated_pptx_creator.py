"""
AnimatedPowerPointCreator — build animated .pptx files programmatically.

Wraps python-pptx and injects the raw Open XML needed for slide transitions
and shape-level entrance / emphasis / exit animations that python-pptx does
not natively support.
"""

from __future__ import annotations

import copy
import uuid
from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional

from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# ── OOX namespaces used when writing raw XML ──────────────────────────
NSMAP = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
}

# ── Transition presets (p:transition child element tag) ────────────────
TRANSITION_MAP = {
    "fade": "p:fade",
    "push": "p:push",
    "wipe": "p:wipe",
    "split": "p:split",
    "cover": "p:cover",
    "cut": "p:cut",
    "dissolve": "p:dissolve",
    "random_bar": "p:randomBar",
}

# ── Animation presets (presetClass / presetID pairs) ───────────────────
ANIMATION_PRESETS = {
    "appear": ("entr", 1),
    "fade_in": ("entr", 10),
    "fly_in_left": ("entr", 2),
    "fly_in_bottom": ("entr", 2),
    "wipe": ("entr", 22),
    "zoom_in": ("entr", 53),
    "bounce": ("entr", 26),
    "fade_out": ("exit", 10),
    "fly_out_right": ("exit", 2),
    "pulse": ("emph", 26),
}


@dataclass
class SlideSpec:
    """Declarative description of a single slide."""

    layout: str = "content"
    title: str = ""
    subtitle: str = ""
    bullets: list[str] = field(default_factory=list)
    image_path: Optional[str] = None
    animation: Optional[str] = None
    stagger_delay: float = 0.4
    notes: str = ""


class AnimatedPowerPointCreator:
    """High-level builder for animated PowerPoint presentations."""

    def __init__(self, template: Optional[str | Path] = None) -> None:
        if template:
            self._prs = Presentation(str(template))
        else:
            self._prs = Presentation()
        self._slide_specs: list[SlideSpec] = []
        self._transitions: dict[int, tuple[str, float]] = {}

    # ── public API ─────────────────────────────────────────────────────

    def add_slide(self, **kwargs) -> int:
        """Add a slide and return its 0-based index."""
        spec = SlideSpec(**kwargs)
        self._slide_specs.append(spec)
        return len(self._slide_specs) - 1

    def set_transition(
        self,
        slide_index: int,
        transition: str = "fade",
        duration: float = 0.5,
    ) -> None:
        if transition not in TRANSITION_MAP:
            raise ValueError(
                f"Unknown transition '{transition}'. "
                f"Choose from: {', '.join(TRANSITION_MAP)}"
            )
        self._transitions[slide_index] = (transition, duration)

    def save(self, path: str | Path) -> Path:
        """Build all slides, apply animations / transitions, and save."""
        path = Path(path)
        for spec in self._slide_specs:
            self._build_slide(spec)

        for idx, (trans, dur) in self._transitions.items():
            if idx < len(self._prs.slides):
                self._apply_transition(self._prs.slides[idx], trans, dur)

        self._prs.save(str(path))
        return path

    # ── internal: slide construction ───────────────────────────────────

    def _pick_layout(self, name: str):
        name_lower = name.lower()
        layout_map = {
            "title": 0,
            "content": 1,
            "section": 2,
            "two_content": 3,
            "blank": 6,
        }
        idx = layout_map.get(name_lower, 1)
        idx = min(idx, len(self._prs.slide_layouts) - 1)
        return self._prs.slide_layouts[idx]

    def _build_slide(self, spec: SlideSpec):
        layout = self._pick_layout(spec.layout)
        slide = self._prs.slides.add_slide(layout)

        if spec.title and slide.placeholders:
            slide.placeholders[0].text = spec.title

        if spec.subtitle and len(slide.placeholders) > 1:
            slide.placeholders[1].text = spec.subtitle

        if spec.bullets and len(slide.placeholders) > 1:
            tf = slide.placeholders[1].text_frame
            tf.clear()
            for i, bullet in enumerate(spec.bullets):
                para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                para.text = bullet
                para.level = 0
                para.font.size = Pt(18)

        if spec.image_path:
            slide.shapes.add_picture(
                spec.image_path,
                Inches(1),
                Inches(2),
                width=Inches(4),
            )

        if spec.notes:
            slide.notes_slide.notes_text_frame.text = spec.notes

        if spec.animation:
            self._apply_animation(slide, spec.animation, spec.stagger_delay)

    # ── internal: transitions (raw XML) ────────────────────────────────

    @staticmethod
    def _apply_transition(slide, transition: str, duration: float):
        tag = TRANSITION_MAP[transition]
        ns_tag = tag.replace("p:", f"{{{NSMAP['p']}}}")

        dur_ms = str(int(duration * 1000))
        trans_el = etree.SubElement(
            slide._element,
            f"{{{NSMAP['p']}}}transition",
            attrib={"spd": "med", "advTm": dur_ms},
        )
        etree.SubElement(trans_el, ns_tag)

    # ── internal: shape animations (raw XML) ───────────────────────────

    @staticmethod
    def _apply_animation(slide, animation: str, stagger: float):
        if animation not in ANIMATION_PRESETS:
            raise ValueError(
                f"Unknown animation '{animation}'. "
                f"Choose from: {', '.join(ANIMATION_PRESETS)}"
            )
        preset_class, preset_id = ANIMATION_PRESETS[animation]

        timing = etree.SubElement(
            slide._element,
            f"{{{NSMAP['p']}}}timing",
        )
        tn_lst = etree.SubElement(timing, f"{{{NSMAP['p']}}}tnLst")
        par = etree.SubElement(tn_lst, f"{{{NSMAP['p']}}}par")
        c_tn = etree.SubElement(
            par,
            f"{{{NSMAP['p']}}}cTn",
            attrib={"id": "1", "dur": "indefinite", "restart": "never", "nodeType": "tmRoot"},
        )
        child_tn_lst = etree.SubElement(c_tn, f"{{{NSMAP['p']}}}childTnLst")

        shapes = [s for s in slide.shapes if s.has_text_frame]
        for i, shape in enumerate(shapes):
            delay_ms = str(int(i * stagger * 1000))
            sp_id = str(shape.shape_id)

            seq_par = etree.SubElement(child_tn_lst, f"{{{NSMAP['p']}}}par")
            seq_ctn = etree.SubElement(
                seq_par,
                f"{{{NSMAP['p']}}}cTn",
                attrib={"id": str(i + 2), "fill": "hold"},
            )
            s_cond = etree.SubElement(seq_ctn, f"{{{NSMAP['p']}}}stCondLst")
            etree.SubElement(
                s_cond,
                f"{{{NSMAP['p']}}}cond",
                attrib={"delay": delay_ms},
            )

            inner_child = etree.SubElement(seq_ctn, f"{{{NSMAP['p']}}}childTnLst")
            inner_par = etree.SubElement(inner_child, f"{{{NSMAP['p']}}}par")
            inner_ctn = etree.SubElement(
                inner_par,
                f"{{{NSMAP['p']}}}cTn",
                attrib={
                    "id": str(100 + i),
                    "presetID": str(preset_id),
                    "presetClass": preset_class,
                    "presetSubtype": "0",
                    "fill": "hold",
                    "nodeType": "afterEffect",
                },
            )
            inner_s_cond = etree.SubElement(inner_ctn, f"{{{NSMAP['p']}}}stCondLst")
            etree.SubElement(
                inner_s_cond,
                f"{{{NSMAP['p']}}}cond",
                attrib={"delay": "0"},
            )


def main():
    """Quick demo — generates a sample animated presentation."""
    creator = AnimatedPowerPointCreator()

    creator.add_slide(
        layout="title",
        title="AnimatedPowerPointCreator",
        subtitle="Build animated .pptx files with Python",
        animation="fade_in",
    )

    creator.add_slide(
        layout="content",
        title="Features",
        bullets=[
            "Slide transitions (fade, push, wipe …)",
            "Shape-level entrance & exit animations",
            "Staggered bullet animations",
            "Template support",
        ],
        animation="fly_in_left",
        stagger_delay=0.3,
    )

    creator.add_slide(
        layout="content",
        title="Getting Started",
        bullets=[
            "pip install -r requirements.txt",
            "Import AnimatedPowerPointCreator",
            "Add slides, set transitions, save!",
        ],
        animation="zoom_in",
    )

    creator.set_transition(1, "fade", 0.7)
    creator.set_transition(2, "push", 0.5)

    out = creator.save("demo_output.pptx")
    print(f"Saved → {out}")


if __name__ == "__main__":
    main()
